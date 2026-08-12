#!/usr/bin/env python3
"""
ingest_v2.py — 本地 RAG 灌庫（P1 重構，2026-06-08）

路線 B hybrid：dense = intfloat/multilingual-e5-large (1024) + sparse = Qdrant/bm25(IDF)
特性：增量(sha256 去重)、按 source_path 刪舊向量、metadata 分類、🔴 secrets 排除、.txt/.md/.json/.yaml/.pdf
設計書：RAG_ENHANCEMENT_DESIGN.md

SMB 或本機 vault 掛好後設 kb_root 即可跑。
"""
import os, re, sys, json, hashlib, datetime, uuid
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter
from qdrant_client import QdrantClient, models
from .config import cfg
from . import embedding, store

# ==================== 設定 ====================
# vault 下的原始檔庫；實際位置由 kb_config.yaml 的 kb_root 控制。
KB_ROOT      = cfg("kb_root")
SOURCE_ROOT  = os.path.join(KB_ROOT, "raw_files")
STATE_PATH   = os.path.join(KB_ROOT, "state", "ingest_manifest.json")

QDRANT_HOST  = cfg("qdrant.host")
QDRANT_PORT  = cfg("qdrant.port")
COLLECTION   = cfg("qdrant.collection")
QDRANT_TIMEOUT = cfg("qdrant.timeout_batch")

DENSE_MODEL  = cfg("embedding.dense_model")   # 1024 維
SPARSE_MODEL = cfg("embedding.sparse_model")
DENSE_DIM    = cfg("embedding.dense_dim")

CHUNK_SIZE, CHUNK_OVERLAP = cfg("chunking.size"), cfg("chunking.overlap")
SUPPORTED_EXT = {
    ".txt", ".md", ".json", ".yaml", ".yml", ".pdf",
    # Office formats: only the modern zip/XML ones. The legacy binary .doc,
    # .xls and .ppt need a different parser entirely and are left out rather
    # than half-supported.
    ".docx", ".xlsx", ".pptx", ".csv",
}


# ==================== 來源（vault 內 + 用戶指定路徑）====================
class Source:
    """One place to read documents from.

    ``raw_files`` inside the vault is always a source and is the one this tool
    writes to. Everything else is a folder the user pointed at — read strictly
    read-only, including external drives that may not be mounted right now.
    """

    def __init__(self, root, label=None, builtin=False):
        self.root = Path(os.path.expanduser(str(root)))
        self.builtin = builtin
        self.label = label or self.root.name or "_uncategorized"

    @property
    def available(self):
        try:
            return self.root.is_dir()
        except OSError:
            # An unmounted volume can raise rather than answer False.
            return False

    def category_of(self, path):
        """Which category a file under this source belongs to.

        raw_files keeps its existing rule — the first folder under it names the
        category — because that is what kb_add writes and what every existing
        vault already looks like. A user-added folder uses its own label for
        everything inside it, so pointing at a drive does not invent a category
        per subfolder.
        """
        if not self.builtin:
            return self.label
        try:
            rel = path.relative_to(self.root)
        except ValueError:
            return "_uncategorized"
        return rel.parts[0] if len(rel.parts) > 1 else "_uncategorized"

    def contains(self, path):
        try:
            Path(path).relative_to(self.root)
            return True
        except ValueError:
            return False


def configured_sources():
    """Every source this deployment reads: raw_files first, then user folders.

    Read at call time rather than import time so a path added in Settings
    applies to the next run without restarting the process.
    """
    from .add import notes_root

    sources = [Source(SOURCE_ROOT, builtin=True)]
    seen = {str(sources[0].root)}

    # Wherever notes are written is always read back, so a custom notes folder
    # still reaches the index and the catalog. Without this, pointing notes
    # somewhere else would quietly stop them being searchable at all.
    notes = Source(notes_root(), builtin=True)
    if str(notes.root) not in seen:
        sources.append(notes)
        seen.add(str(notes.root))

    for entry in cfg("sources", []) or []:
        if isinstance(entry, dict):
            raw, label = entry.get("path"), entry.get("label")
        else:
            raw, label = entry, None
        if not raw:
            continue
        source = Source(raw, label=label)
        key = str(source.root)
        if key in seen:
            continue
        seen.add(key)
        sources.append(source)
    return sources


def max_file_bytes():
    """Cap on a single file, so pointing at a drive cannot pull in a 2GB blob."""
    try:
        return max(1, int(cfg("ingest.max_file_mb", 25))) * 1024 * 1024
    except (TypeError, ValueError):
        return 25 * 1024 * 1024

# ==================== 🔴 Secrets 排除（硬規）====================
# 硬名單：檔名本身就係憑證載體 → 無條件跳過
SECRET_NAME_HARD = re.compile(r"(^\.env|\.key$|^id_rsa|\.pem$|\.p12$)", re.I)
# 軟名單：檔名提及憑證字眼（可能只係「講 token 嘅筆記」）→ 內容都命中先跳過
SECRET_NAME_SOFT = re.compile(r"(^|[._-])(secret|credential|password|token|auth)", re.I)
# 明確白名單：符合以下模式的檔名即使觸發軟名單也不跳過
SECRET_NAME_WHITELIST = re.compile(r"design[-_]tokens?", re.I)
SECRET_CONTENT_PATS = [
    re.compile(r"sk-[A-Za-z0-9]{20,}"),                 # OpenAI/DeepSeek 類
    re.compile(r"gh[pousr]_[A-Za-z0-9]{20,}"),          # GitHub token
    re.compile(r"xox[baprs]-[A-Za-z0-9-]{10,}"),        # Slack
    re.compile(r"AKIA[0-9A-Z]{16}"),                    # AWS
    re.compile(r"(api[_-]?key|secret|password)\s*[:=]\s*['\"]?[A-Za-z0-9/+_\-]{16,}", re.I),
]

def looks_secret(path: Path, text: str) -> bool:
    if SECRET_NAME_HARD.search(path.name):
        return True                                   # .env/.key/id_rsa 類：名就係證據
    head = text[:20000]
    content_hit = any(p.search(head) for p in SECRET_CONTENT_PATS)
    if SECRET_NAME_WHITELIST.search(path.name):
        return content_hit                            # 白名單檔名照查內容
    if SECRET_NAME_SOFT.search(path.name):
        return content_hit                            # 「講憑證嘅筆記」內容乾淨就放行
    return content_hit

# ==================== 敏感度（初期極簡：public/sensitive）====================
SENSITIVE_HINT = re.compile(r"(機密|身分證|身份證|個人資料|個資|學生名單|成績單|輔導紀錄|病歷|信用卡|confidential|\bprivate\b|\bsecret\b)", re.I)
def classify_sensitivity(path: Path, text: str) -> str:
    if SENSITIVE_HINT.search(str(path)) or SENSITIVE_HINT.search(text[:5000]):
        return "sensitive"
    return "public"

# ==================== 讀檔 ====================
# A single document cannot contribute more text than this. A 5MB spreadsheet
# can hold millions of cells; without a ceiling one file could dominate both
# the run time and the index.
MAX_TEXT_CHARS = 2_000_000


def _load_docx(path: Path) -> str:
    """Paragraphs and table cells, in document order where the API allows."""
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _load_xlsx(path: Path) -> str:
    """Cell values sheet by sheet. Formulas read as their cached values."""
    from openpyxl import load_workbook

    book = load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts, total = [], 0
        for sheet in book.worksheets:
            parts.append(f"# {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if not cells:
                    continue
                line = "\t".join(cells)
                total += len(line)
                parts.append(line)
                if total > MAX_TEXT_CHARS:
                    parts.append("…（內容過長，已截斷）")
                    return "\n".join(parts)
        return "\n".join(parts)
    finally:
        book.close()


def _load_pptx(path: Path) -> str:
    """Text frames and table cells, slide by slide."""
    from pptx import Presentation

    deck = Presentation(str(path))
    parts = []
    for number, slide in enumerate(deck.slides, 1):
        parts.append(f"# Slide {number}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append("\t".join(cells))
    return "\n".join(parts)


def _load_csv(path: Path) -> str:
    """Rows as tab-separated lines, sniffing the delimiter where possible."""
    import csv

    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        dialect = csv.Sniffer().sniff(raw[:4096])
    except csv.Error:
        dialect = csv.excel
    lines = []
    for row in csv.reader(raw.splitlines(), dialect):
        cells = [c.strip() for c in row if c and c.strip()]
        if cells:
            lines.append("\t".join(cells))
    return "\n".join(lines)


def _load_pdf(path: Path) -> str:
    from langchain_community.document_loaders import PyPDFLoader

    return "\n".join(d.page_content for d in PyPDFLoader(str(path)).load())


# Office formats are read through their own libraries rather than a generic
# extractor: each is a zip of XML, and the alternative (unstructured) pulls in
# a far larger dependency tree than a desktop bundle should carry.
LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".xlsx": _load_xlsx,
    ".pptx": _load_pptx,
    ".csv": _load_csv,
}
PLAIN_EXT = {".txt", ".md", ".json", ".yaml", ".yml"}


def load_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in PLAIN_EXT:
        text = path.read_text(encoding="utf-8", errors="replace")
    else:
        loader = LOADERS.get(ext)
        text = loader(path) if loader else ""
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + "\n…（內容過長，已截斷）"
    return text

def sha256(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for b in iter(lambda: f.read(65536), b""):
            h.update(b)
    return h.hexdigest()

def chunk_id(source_path: str, idx: int) -> str:
    return str(uuid.uuid5(uuid.NAMESPACE_URL, f"{source_path}::{idx}"))

# ==================== Qdrant collection ====================
def ensure_collection(client: QdrantClient):
    if client.collection_exists(COLLECTION):
        return
    client.create_collection(
        collection_name=COLLECTION,
        vectors_config={"dense": models.VectorParams(size=DENSE_DIM, distance=models.Distance.COSINE)},
        sparse_vectors_config={"sparse": models.SparseVectorParams(modifier=models.Modifier.IDF)},
    )
    print(f"✅ 已建 collection {COLLECTION}（dense {DENSE_DIM} + sparse bm25/IDF）")

# ==================== 主流程 ====================
def main():
    sources = configured_sources()
    available = [s for s in sources if s.available]
    offline = [s for s in sources if not s.available]
    if not available:
        sys.exit(f"❌ 冇任何來源可讀（vault 原始檔庫：{SOURCE_ROOT}）")
    for s in offline:
        # Not fatal, and deliberately loud: the index for this source is kept
        # rather than pruned, so the user has to be told why nothing changed.
        print(f"  ⏸️  來源唔喺度，今次跳過（已入庫嘅內容保留）：{s.root}")

    client = store.connect(QDRANT_TIMEOUT)
    ensure_collection(client)

    # 載入 manifest（增量基準）
    manifest = {}
    if os.path.exists(STATE_PATH):
        manifest = json.loads(Path(STATE_PATH).read_text(encoding="utf-8"))

    print("🔎 載入 embedding 模型（首次會下載 e5-large ~2.24GB）...")
    dense, sparse = embedding.pair(DENSE_MODEL, SPARSE_MODEL)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", "。", "！", "？", " ", ""],
    )

    seen_paths, to_delete_ids, new_points = set(), [], []
    skipped_secret, unchanged, processed = [], 0, 0

    size_cap = max_file_bytes()
    skipped_large = []

    for source in available:
        print(f"📂 掃描來源：{source.root}（類別："
              f"{'依子資料夾' if source.builtin else source.label}）")
        for fp in source.root.rglob("*"):
            try:
                if fp.is_dir() or fp.name.startswith("."):
                    continue
            except OSError:
                # A volume yanked mid-scan; stop trusting this source's tree.
                print(f"  ⚠️  來源讀唔到，中止掃描：{source.root}")
                break
            if fp.suffix.lower() not in SUPPORTED_EXT:
                continue
            try:
                if fp.stat().st_size > size_cap:
                    skipped_large.append(str(fp))
                    print(f"  ⏭️  太大，跳過（>{size_cap // (1024*1024)}MB）：{fp.name}")
                    continue
            except OSError:
                continue
            spath = str(fp)
            seen_paths.add(spath)
            try:
                digest = sha256(fp)
            except OSError as e:
                print(f"  ⚠️  讀唔到，跳過：{fp.name}（{type(e).__name__}）")
                continue

            # 增量：未變 → skip
            if manifest.get(spath, {}).get("hash") == digest:
                unchanged += 1
                continue

            try:
                text = load_text(fp)
            except Exception as e:
                # Broad on purpose: a corrupt .docx, a password-protected
                # .xlsx or a missing optional parser must cost one file, not
                # the whole run.
                print(f"  ⚠️  讀唔到，跳過：{fp.name}（{type(e).__name__}: {e}）")
                continue
            if looks_secret(fp, text):
                skipped_secret.append(spath)
                print(f"  🔴 SKIP（疑似 secrets，不入庫）：{fp.name}")
                continue

            # 變動檔：先刪舊向量
            if spath in manifest:
                to_delete_ids += manifest[spath].get("chunk_ids", [])

            category = source.category_of(fp)
            sensitivity = classify_sensitivity(fp, text)
            now = datetime.date.today().isoformat()
            chunks = [c.page_content for c in splitter.create_documents([text])]
            if not chunks:
                manifest[spath] = {"hash": digest, "chunk_ids": [], "ingested_at": now}
                continue

            dvecs = list(dense.embed([f"passage: {c}" for c in chunks]))  # e5 需 passage: 前綴
            svecs = list(sparse.embed(chunks))
            ids = []
            for i, (c, dv, sv) in enumerate(zip(chunks, dvecs, svecs)):
                pid = chunk_id(spath, i)
                ids.append(pid)
                new_points.append(models.PointStruct(
                    id=pid,
                    vector={
                        "dense": dv.tolist(),
                        "sparse": models.SparseVector(indices=sv.indices.tolist(), values=sv.values.tolist()),
                    },
                    payload={
                        "text": c, "source_path": spath, "source_file": fp.name,
                        "category": category, "file_type": fp.suffix.lower().lstrip("."),
                        "sensitivity": sensitivity, "file_hash": digest,
                        "chunk_index": i, "ingested_at": now,
                    },
                ))
            manifest[spath] = {"hash": digest, "chunk_ids": ids, "ingested_at": now}
            processed += 1

    # 已從磁碟消失的檔 → 刪向量 + 出 manifest。
    #
    # 唯一可以刪嘅前提：嗰個檔所屬嘅來源今次真係掃過。一個未掛載嘅外置碟，
    # 佢下面每個檔都會「唔見咗」，若照刪就等於拔一次 USB 就殺晒成個索引。
    # 所以未掛載嘅來源一律保留；只有真係掃過（檔案確認消失），或者用戶已經
    # 喺設定度移走咗嗰個來源，先至清走。
    kept_offline = 0
    for gone in [p for p in manifest if p not in seen_paths]:
        owner = next((s for s in sources if s.contains(gone)), None)
        if owner is not None and not owner.available:
            kept_offline += 1
            continue
        to_delete_ids += manifest[gone].get("chunk_ids", [])
        del manifest[gone]
        print(f"  🗑️  來源已刪，移除向量：{os.path.basename(gone)}")
    if kept_offline:
        print(f"  💾 保留 {kept_offline} 個未掛載來源嘅已入庫檔案（冇刪向量）")

    if to_delete_ids:
        client.delete(collection_name=COLLECTION,
                      points_selector=models.PointIdsList(points=to_delete_ids))
        print(f"🗑️  刪除舊/孤兒向量 {len(to_delete_ids)} 點")

    for i in range(0, len(new_points), 128):
        client.upsert(collection_name=COLLECTION, wait=True, points=new_points[i:i+128])
    if new_points:
        print(f"🧠 upsert 新向量 {len(new_points)} 點")

    Path(STATE_PATH).parent.mkdir(parents=True, exist_ok=True)
    Path(STATE_PATH).write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"\n🎉 完成：處理 {processed} 檔 / 未變 {unchanged} 檔 / 跳過 secrets {len(skipped_secret)} 檔"
          + (f" / 太大 {len(skipped_large)} 檔" if skipped_large else ""))
    print(f"   manifest → {STATE_PATH}")

    # Loop 1：灌庫後自動更新知識目錄（增量；失敗唔阻斷 ingest）
    if processed or to_delete_ids:
        try:
            from . import index_update
            print("\n📚 觸發知識目錄更新（Loop 1）…")
            index_update.main()
        except Exception as e:
            print(f"⚠️ 目錄更新略過（不影響灌庫）：{e}")

    # 更新 orctl KB 快取 → /tmp/orctl_kb_cache.json
    # ingest 在終端/有 TCC 環境跑；launchd 的 orctl-web 讀此快取繞過 TCC 限制
    try:
        import subprocess as _sp, sys as _sys
        _sp.run([_sys.executable, str(Path.home() / "bin/orctl"), "kb", "--json"],
                capture_output=True, timeout=10)
    except Exception:
        pass

if __name__ == "__main__":
    main()
